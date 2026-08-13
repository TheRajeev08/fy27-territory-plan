"""Geometry and consistency checks for a generated deck.

    python3 verify_deck.py <deck.pptx> [--coverage coverage.json]

Three classes of defect, all of which have actually shipped in this deck before:

  * Overflow - a shape whose text needs more height than the shape has, or whose box
    runs off the slide. Renders as clipped text on someone else's screen, not yours.
  * Collision - two body shapes overlapping. Usually a table that grew a row.
  * Inconsistency - a number on one slide contradicting the same number on another,
    or contradicting coverage.json. This is the one that costs credibility in the room.

Exits non-zero when anything is found, so it can gate a release.
"""
import json
import os
import re
import sys

from pptx import Presentation
from pptx.util import Emu

# Rough advance width per point of font size, averaged over mixed-case text. Only used
# to estimate wrapping, so it is deliberately conservative: better to warn about a box
# that fits than to pass one that clips.
CHAR_W = 0.50
LINE_H = 1.22

# Advance width used for the *width* check. Deliberately smaller than CHAR_W: over-
# estimating height only costs a spurious warning, but over-estimating width would
# flag boxes that render fine, and a verifier that cries wolf stops being read.
CHAR_W_TIGHT = 0.42
WIDTH_TOLERANCE = 1.02


def shape_text(shape):
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def font_size(shape, default=12.0):
    sizes = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return max(sizes) if sizes else default


def letter_spacing(shape):
    """Extra advance per character, in points, from the `spc` run property."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            spc = run.font._rPr.get("spc")
            if spc:
                return float(spc) / 100.0
    return 0.0


def needed_width(shape):
    """Estimated rendered width in EMU for text that will not wrap.

    Only meaningful when word_wrap is False. Such text renders on one line and spills
    sideways out of its box, so a height check can never catch it - which is exactly
    how 13 clipped cells shipped on slide 2. Returns 0 when the check does not apply.
    """
    if not shape.has_text_frame or shape.text_frame.word_wrap is not False:
        return 0
    text = shape_text(shape)
    if not text.strip():
        return 0
    size = font_size(shape)
    spacing = letter_spacing(shape)
    widest = 0.0
    for line in text.split("\n"):
        widest = max(widest, len(line) * (size * CHAR_W_TIGHT + spacing))
    return int(widest / 72.0 * 914400)


def needed_height(shape):
    """Estimated rendered text height in EMU."""
    text = shape_text(shape)
    if not text.strip():
        return 0
    size = font_size(shape)
    width_in = Emu(shape.width).inches
    if width_in <= 0:
        return 0
    # word_wrap=False renders on a single line regardless of width: it spills sideways
    # rather than growing downwards, so estimating wrapped height would be wrong.
    wrap = shape.text_frame.word_wrap
    if wrap is False:
        lines = len(text.split("\n"))
    else:
        chars_per_line = max(1, int(width_in * 72.0 / (size * CHAR_W)))
        lines = 0
        for para in text.split("\n"):
            lines += max(1, -(-len(para) // chars_per_line))
    return int(lines * size * LINE_H / 72.0 * 914400)


def boxes_overlap(a, b):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ix = min(ax2, bx2) - max(ax1, bx1)
    iy = min(ay2, by2) - max(ay1, by1)
    if ix <= 0 or iy <= 0:
        return 0.0
    smaller = min((ax2 - ax1) * (ay2 - ay1), (bx2 - bx1) * (by2 - by1))
    return (ix * iy) / smaller if smaller else 0.0


def check_geometry(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    overflows, collisions = [], []

    for index, slide in enumerate(prs.slides, 1):
        boxes = []
        for shape in slide.shapes:
            if shape.width is None or shape.height is None:
                continue
            box = (shape.left, shape.top, shape.left + shape.width,
                   shape.top + shape.height)

            if box[2] > sw + 1000 or box[3] > sh + 1000 or box[0] < -1000 or box[1] < -1000:
                overflows.append((index, "off-slide", shape.shape_type,
                                  shape_text(shape)[:60]))
            need = needed_height(shape)
            if need and need > shape.height * 1.06:
                overflows.append((index, "text %.2fin > box %.2fin"
                                  % (Emu(need).inches, Emu(shape.height).inches),
                                  shape.shape_type, shape_text(shape)[:60]))
            wide = needed_width(shape)
            if wide and wide > shape.width * WIDTH_TOLERANCE:
                overflows.append((index, "clipped %.2fin wide > box %.2fin"
                                  % (Emu(wide).inches, Emu(shape.width).inches),
                                  shape.shape_type, shape_text(shape)[:60]))
            if shape.has_text_frame and shape_text(shape).strip():
                boxes.append((box, shape_text(shape)[:40]))

        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                frac = boxes_overlap(boxes[i][0], boxes[j][0])
                if frac > 0.30:
                    collisions.append((index, round(frac, 2), boxes[i][1], boxes[j][1]))

    return len(prs.slides), overflows, collisions


MONEY = re.compile(r"\$[\d,]+(?:\.\d+)?[KM]?")


def deck_money(value):
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from deck import money
    return money(value)


def money_tokens(path):
    prs = Presentation(path)
    found = {}
    for index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            for token in MONEY.findall(shape_text(shape)):
                found.setdefault(token, set()).add(index)
    return found


def check_consistency(path, coverage):
    """Assert the headline numbers on the slides match coverage.json."""
    problems = []
    tokens = money_tokens(path)

    def fmt(value):
        return deck_money(value)

    totals = coverage.get("totals", {})
    pipeline = coverage.get("pipeline", {})
    by_bucket = pipeline.get("byBucket", {}) or {}

    expected = {
        "Bucket 1 live pipeline": by_bucket.get("Bucket 1"),
        "Bucket 2 live pipeline": by_bucket.get("Bucket 2"),
        "H1 renewal pipeline": pipeline.get("renewal"),
    }
    for label, value in expected.items():
        if not value:
            continue
        if fmt(value) not in tokens:
            problems.append("%s %s does not appear on any slide" % (label, fmt(value)))

    # The blended net-new figure is exactly the number this change set removed. If it
    # reappears, a slide has regressed to presenting Bucket 2 money as Bucket 1 cover.
    blended = pipeline.get("netNew")
    if blended and fmt(blended) in tokens:
        problems.append(
            "blended net-new %s appears on slide(s) %s - it mixes buckets and must not "
            "be shown beside Bucket 1 coverage"
            % (fmt(blended), sorted(tokens[fmt(blended)])))

    if totals.get("targetsComplete"):
        target, attained, gap = (totals.get("h1Target"), totals.get("attainedH1"),
                                 totals.get("gap"))
        if round(target - attained, 2) != round(gap, 2):
            problems.append("totals do not reconcile: %s - %s != %s"
                            % (target, attained, gap))
    return problems


def main():
    if len(sys.argv) < 2:
        raise SystemExit(__doc__)
    path = sys.argv[1]
    slides, overflows, collisions = check_geometry(path)

    coverage_path = (sys.argv[sys.argv.index("--coverage") + 1]
                     if "--coverage" in sys.argv else "")
    problems = []
    if coverage_path and os.path.exists(coverage_path):
        with open(coverage_path, "r", encoding="utf-8") as fh:
            problems = check_consistency(path, json.load(fh))

    print(json.dumps({
        "deck": path,
        "slides": slides,
        "overflows": len(overflows),
        "collisions": len(collisions),
        "inconsistencies": len(problems),
    }))
    for slide, why, kind, text in overflows:
        print("  overflow  slide %d: %s [%s] %r" % (slide, why, kind, text))
    for slide, frac, a, b in collisions:
        print("  collision slide %d: %.0f%% overlap %r / %r" % (slide, frac * 100, a, b))
    for problem in problems:
        print("  mismatch  %s" % problem)

    if overflows or collisions or problems:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
