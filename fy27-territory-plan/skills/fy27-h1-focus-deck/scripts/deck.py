"""Render the FY27 H1 focus-account leadership deck.

Every number on every slide is read from the pipeline outputs. Nothing here is
composed by hand or by a model at render time, because the whole point of the deck is
that a leader can challenge any figure and it traces back to the workbook, to Kusto, or
to a cited source. If a value is missing, the slide says so rather than estimating.

    deck.py <report.json> <potential.json> <focus-accounts.json> <runDir> \
            [--paf paf.json] [--partners partners.json] [--out deck.pptx]
"""

import json
import os
import sys
from datetime import date

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Emu, Inches, Pt
except ImportError:
    sys.stderr.write(
        "This step needs the 'python-pptx' package to build the presentation.\n"
        "Install it with:  python3 -m pip install --user python-pptx\n"
        "Then run the skill again - nothing else needs to be redone.\n")
    sys.exit(3)

# GitHub-ish dark palette. Dark decks read better projected in a review room, and the
# accent colours are reused consistently so a play is recognisable by colour alone.
INK = RGBColor(0x0D, 0x11, 0x17)
PANEL = RGBColor(0x16, 0x1B, 0x22)
PANEL_2 = RGBColor(0x1C, 0x21, 0x28)
LINE = RGBColor(0x30, 0x36, 0x3D)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x94, 0x9E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PLAY_COLOR = {
    "Innovate": RGBColor(0x58, 0xA6, 0xFF),
    "Trust": RGBColor(0x3F, 0xB9, 0x50),
    "Scale": RGBColor(0xBC, 0x8C, 0xFF),
}
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
WARN = RGBColor(0xD2, 0x99, 0x22)
GOOD = RGBColor(0x3F, 0xB9, 0x50)

# `.title()` mangles acronyms into "Ghec"/"Copilot Aiu", which reads as a typo on a
# leadership slide. Billing product names are lowercase in Kusto, so map them explicitly.
PRODUCT_LABEL = {
    "copilot": "Copilot", "copilot aiu": "Copilot AI credits", "actions": "Actions",
    "ghec": "GHEC", "ghas": "GHAS", "ghcs": "GHCS", "ghsp": "GHSP",
    "ghe": "GHE", "ghes": "GHES", "codespaces": "Codespaces", "packages": "Packages",
    "git_lfs": "Git LFS", "shared_storage": "Shared storage",
}

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)
BODY_TOP = Inches(1.62)


# ---------------------------------------------------------------- formatting

def money(value):
    """Compact currency. Leadership slides read better in $1.2M than $1,234,567."""
    value = float(value or 0)
    if abs(value) >= 1_000_000:
        return "$%.2fM" % (value / 1_000_000)
    if abs(value) >= 1_000:
        return "$%.0fK" % (value / 1_000)
    return "$%.0f" % value


def num(value):
    return "{:,.0f}".format(float(value or 0))


def pct(part, whole):
    return "0%" if not whole else "%.0f%%" % (100.0 * part / whole)


def truncate(text, limit):
    text = (text or "").strip()
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "\u2026"


# Tier labels are owned by rank.py and carry a descriptive suffix ("Tier 1 - Must win").
# The deck matches on the "Tier N" prefix so renaming a tier there does not break rendering.
def tier_names(focus):
    """Tier labels in rank order, taken from the data rather than assumed."""
    seen = []
    for account in focus.get("accounts", []):
        tier = account.get("tier")
        if tier and tier not in seen:
            seen.append(tier)
    return seen


def tier_count(focus, prefix):
    return sum(1 for a in focus.get("accounts", [])
               if (a.get("tier") or "").startswith(prefix))


def tier_label(focus, prefix, fallback):
    for tier in tier_names(focus):
        if tier.startswith(prefix):
            return tier
    return fallback


# ---------------------------------------------------------------- primitives

class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]
        self.notes = []

    def slide(self, title=None, eyebrow=None, note=None):
        slide = self.prs.slides.add_slide(self.blank)
        self.fill(slide, 0, 0, W, H, INK)
        if title:
            if eyebrow:
                self.text(slide, MARGIN, Inches(0.44), W - 2 * MARGIN, Inches(0.26),
                          eyebrow.upper(), size=11, color=ACCENT, bold=True, space=True)
            self.text(slide, MARGIN, Inches(0.72), W - 2 * MARGIN, Inches(0.52),
                      title, size=27, color=WHITE, bold=True)
            self.fill(slide, MARGIN, Inches(1.36), Inches(1.5), Inches(0.035), ACCENT)
        if note:
            self.notes.append((slide, note))
        return slide

    def fill(self, slide, x, y, w, h, color, line=None, radius=False):
        from pptx.enum.shapes import MSO_SHAPE
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, int(x), int(y), int(w), int(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
            shape.line.width = Pt(0.75)
        shape.shadow.inherit = False
        if shape.has_text_frame:
            shape.text_frame.clear()
        return shape

    def text(self, slide, x, y, w, h, value, size=13, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, space=False, anchor=MSO_ANCHOR.TOP, wrap=True):
        box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
        frame = box.text_frame
        frame.word_wrap = wrap
        frame.vertical_anchor = anchor
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        para = frame.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = value
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
        if space:
            # Letter-spacing has no python-pptx API; `spc` is a plain (un-namespaced)
            # attribute on the run properties element, in hundredths of a point.
            run.font._rPr.set("spc", "120")
        return box

    def bullets(self, slide, x, y, w, items, size=13, color=TEXT, gap=0.16,
                bullet="\u25aa", bullet_color=None):
        """Manual bullets with measured wrapping.

        python-pptx cannot measure text, so line count is estimated from the average
        glyph width of the font. Advancing by measured lines rather than a flat step
        keeps long and short bullets evenly spaced instead of leaving ragged gaps.
        """
        width_pt = float(w - Inches(0.24)) / 12700.0
        chars_per_line = max(12, int(width_pt / (size * 0.5)))
        line_h = Inches(size * 1.3 / 72.0)
        cursor = float(y)
        for item in items:
            lines = max(1, -(-len(item) // chars_per_line))
            self.text(slide, x, Emu(int(cursor)), Inches(0.18), line_h, bullet,
                      size=size - 2, color=bullet_color or ACCENT, bold=True)
            self.text(slide, Emu(int(x + Inches(0.24))), Emu(int(cursor)),
                      Emu(int(w - Inches(0.24))), Emu(int(line_h * lines)), item,
                      size=size, color=color)
            cursor += line_h * lines + Inches(gap)
        return Emu(int(cursor))

    def card(self, slide, x, y, w, h, label, value, sub=None, color=None):
        self.fill(slide, x, y, w, h, PANEL, line=LINE, radius=True)
        self.text(slide, Emu(int(x + Inches(0.22))), Emu(int(y + Inches(0.18))),
                  Emu(int(w - Inches(0.4))), Inches(0.24), label.upper(),
                  size=9.5, color=MUTED, bold=True, space=True)
        self.text(slide, Emu(int(x + Inches(0.22))), Emu(int(y + Inches(0.46))),
                  Emu(int(w - Inches(0.4))), Inches(0.5), value,
                  size=25, color=color or WHITE, bold=True)
        if sub:
            self.text(slide, Emu(int(x + Inches(0.22))), Emu(int(y + h - Inches(0.42))),
                      Emu(int(w - Inches(0.4))), Inches(0.3), sub, size=9.5, color=MUTED)

    def table(self, slide, x, y, w, headers, widths, rows, row_h=0.285, size=9.5,
              head_size=9, colors=None):
        """Hand-drawn table. python-pptx native tables fight dark themes."""
        total = float(sum(widths))
        cols = [w * (width / total) for width in widths]
        self.fill(slide, x, y, w, Inches(0.3), PANEL_2)
        cursor_x = float(x)
        for header, col_w in zip(headers, cols):
            self.text(slide, Emu(int(cursor_x + Inches(0.1))), Emu(int(y + Inches(0.055))),
                      Emu(int(col_w - Inches(0.12))), Inches(0.22), header.upper(),
                      size=head_size, color=MUTED, bold=True, space=True, wrap=False)
            cursor_x += col_w
        cursor_y = float(y) + Inches(0.3)
        for index, row in enumerate(rows):
            if index % 2 == 0:
                self.fill(slide, x, Emu(int(cursor_y)), w, Inches(row_h), PANEL)
            cursor_x = float(x)
            row_color = (colors or {}).get(index)
            for cell_index, (cell, col_w) in enumerate(zip(row, cols)):
                color = row_color if (row_color and cell_index == 0) else TEXT
                self.text(slide, Emu(int(cursor_x + Inches(0.1))),
                          Emu(int(cursor_y + Inches(row_h * 0.16))),
                          Emu(int(col_w - Inches(0.12))), Inches(row_h),
                          str(cell), size=size,
                          color=color if cell_index == 0 else color,
                          bold=cell_index == 0, wrap=False)
                cursor_x += col_w
            cursor_y += Inches(row_h)
        return Emu(int(cursor_y))

    def bullet_height(self, w, items, size=13, gap=0.16):
        """Height `bullets` will occupy, so panels can be sized to their content."""
        width_pt = float(w - Inches(0.24)) / 12700.0
        chars_per_line = max(12, int(width_pt / (size * 0.5)))
        line_h = Inches(size * 1.3 / 72.0)
        total = 0.0
        for item in items:
            total += line_h * max(1, -(-len(item) // chars_per_line)) + Inches(gap)
        return total - (Inches(gap) if items else 0)

    def panel(self, slide, x, y, w, label, items, color, size=12, gap=0.18,
              min_h=Inches(1.2)):
        """Titled panel sized to its bullets, so slides never end in dead space."""
        inner_w = w - Inches(0.52)
        body_h = self.bullet_height(inner_w, items, size=size, gap=gap)
        height = max(min_h, body_h + Inches(1.06))
        self.fill(slide, x, y, w, height, PANEL, line=LINE, radius=True)
        self.text(slide, Emu(int(x + Inches(0.26))), Emu(int(y + Inches(0.22))),
                  Emu(int(inner_w)), Inches(0.3), label.upper(), size=10, color=color,
                  bold=True, space=True)
        self.bullets(slide, Emu(int(x + Inches(0.26))), Emu(int(y + Inches(0.66))),
                     Emu(int(inner_w)), items, size=size, gap=gap, bullet_color=color)
        return height

    def footnote(self, slide, value):
        self.text(slide, MARGIN, Inches(6.96), W - 2 * MARGIN, Inches(0.3),
                  value, size=8.5, color=MUTED)

    def save(self, path):
        for slide, note in self.notes:
            slide.notes_slide.notes_text_frame.text = note
        self.prs.save(path)


# ---------------------------------------------------------------- slides

def slide_title(deck, report, focus, potential):
    slide = deck.slide()
    deck.fill(slide, Inches(0.62), Inches(2.42), Inches(0.075), Inches(1.55), ACCENT)
    deck.text(slide, Inches(0.95), Inches(2.42), Inches(11), Inches(0.4),
              "FY27 H1 TERRITORY PLAN", size=13, color=ACCENT, bold=True, space=True)
    deck.text(slide, Inches(0.95), Inches(2.9), Inches(11.4), Inches(0.9),
              "Focus Accounts for the Half", size=44, color=WHITE, bold=True)
    deck.text(slide, Inches(0.95), Inches(3.85), Inches(11), Inches(0.4),
              "%d accounts selected from a %d-account book \u2014 sized, sequenced and evidenced"
              % (focus.get("selectedCount", 0), focus.get("bookSize", 0)),
              size=15, color=MUTED)

    totals = focus.get("totals", {})
    cards = [
        ("Focus accounts", str(focus.get("selectedCount", 0)), "of %d in book" % focus.get("bookSize", 0)),
        ("Potential ARR", money(totals.get("potentialArr")), "new business, sized"),
        ("Current ARR", money(totals.get("currentArr")), "installed in focus set"),
        ("Two-way engaged", str(totals.get("withTwoWay", 0)), "verified conversations"),
    ]
    x = Inches(0.95)
    for label, value, sub in cards:
        deck.card(slide, x, Inches(4.7), Inches(2.75), Inches(1.35), label, value, sub)
        x += Inches(2.92)
    deck.footnote(slide, "Source: %s \u00b7 generated %s \u00b7 all figures traceable to the "
                         "territory workbook and Kusto billing facts."
                  % (report.get("sourceName", "SuperDash export"), date.today().isoformat()))
    return slide


def slide_thesis(deck, report, focus, potential):
    slide = deck.slide("The shape of the half", "Where the half is won",
                       note="Lead with the greenfield reality. This book is not a renewal "
                            "book - the majority of accounts carry no ARR today, so the plan "
                            "is sized on potential and must be judged on pipeline creation.")
    installed = potential.get("installed", {})
    sized = potential.get("accountsSized", 0)
    total = potential.get("accountsTotal", 0)
    with_arr = potential.get("accountsWithArr", 0)

    deck.panel(slide, MARGIN, BODY_TOP, Inches(6.05), "The honest starting point", [
        "%d of %d accounts carry ARR today. This is a greenfield book, so the plan is "
        "sized on potential, not on installed base." % (with_arr, total),
        "%d of %d accounts carry enough product signal to size in dollars (%s). The rest "
        "need discovery before they can be forecast."
        % (sized, total, pct(sized, total)),
        "Installed ARR across the book is %s with %s of annualised consumption revenue "
        "running alongside it." % (money(installed.get("arr")),
                                   money(sum((installed.get("acrAnnualised") or {}).values()))),
        "Potential and installed are never added together anywhere in this deck.",
    ], WARN, size=12.5, gap=0.22)

    deck.panel(slide, Inches(7.0), BODY_TOP, Inches(5.7), "So the half is built this way", [
        "Concentrate on %d accounts rather than spreading across %d. Coverage is not a "
        "strategy in a greenfield book." % (focus.get("selectedCount", 0), focus.get("bookSize", 0)),
        "Rank on three things only: sized potential, verified two-way communication, and "
        "a dated live trigger.",
        "Run one play per account, sequenced from the GitHub Product Adoption Framework "
        "so execution is repeatable rather than improvised.",
        "Tier the list so Tier 1 gets weekly attention and Tier 3 gets a nurture cadence.",
    ], GOOD, size=12.5, gap=0.22)
    deck.footnote(slide, "Sizing coverage and ARR coverage are stated on every slide that "
                         "uses them; no account is assumed to be worth more than its evidence.")
    return slide


def slide_accounts(deck, focus, tier, index, of):
    accounts = [a for a in focus.get("accounts", []) if a.get("tier") == tier]
    if not accounts:
        return None
    label = {"Tier 1": "Weekly cadence \u00b7 executive-sponsored",
             "Tier 2": "Bi-weekly cadence \u00b7 seller-led",
             "Tier 3": "Monthly nurture \u00b7 trigger-activated"}.get(tier[:6], "")
    slide = deck.slide("%s \u2014 %d accounts" % (tier.replace(" - ", " \u00b7 "), len(accounts)),
                       "Q1 \u00b7 Key accounts (%d of %d)" % (index, of),
                       note="Ranked by composite of potential ARR, active communication and "
                            "live trigger recency. 'Why now' is the highest-weighted "
                            "trigger where one was found, otherwise the strongest product signal.")
    deck.text(slide, MARGIN, Inches(1.44), Inches(8), Inches(0.24), label,
              size=10.5, color=MUTED)

    def why_now(account):
        """Prefer a dated trigger; otherwise the evidence line matching the account's play."""
        triggers = account.get("triggers") or []
        if triggers:
            return triggers[0].get("headline", "")
        evidence = account.get("evidence") or []
        play = (account.get("play") or "").lower()
        for line in evidence:
            if play and line.lower().startswith(play):
                return line
        return evidence[0] if evidence else ""

    rows, colors = [], {}
    for i, account in enumerate(accounts[:16]):
        rows.append([
            truncate(account.get("name", ""), 26),
            account.get("play", "\u2014"),
            money(account.get("potentialArr")),
            money((account.get("current") or {}).get("arr")),
            "Yes" if account.get("twoWay") else "\u2014",
            truncate(why_now(account), 68),
        ])
        colors[i] = PLAY_COLOR.get(account.get("play"), TEXT)

    deck.table(slide, MARGIN, Inches(1.78), W - 2 * MARGIN,
               ["Account", "Play", "Potential ARR", "Current ARR", "2-way", "Why now"],
               [3.0, 1.15, 1.35, 1.25, 0.7, 5.0], rows, colors=colors)
    deck.footnote(slide, "Potential ARR is new business only and excludes existing consumption. "
                         "Current ARR from Kusto seat facts; blank means no billing record.")
    return slide


def slide_plays(deck, report, focus, paf):
    slide = deck.slide("One play per account, three plays across the half",
                       "Q2 \u00b7 Plays",
                       note="Play assignment comes from observed product signals in the "
                            "SuperDash export, not from seller preference. Each play has a "
                            "PAF-sequenced set of key actions behind it.")
    mix = focus.get("playMix", {})
    accounts = focus.get("accounts", [])
    x = MARGIN
    width = (W - 2 * MARGIN - Inches(0.5)) / 3
    for play in ("Innovate", "Trust", "Scale"):
        color = PLAY_COLOR.get(play, ACCENT)
        count = mix.get(play, 0)
        pot = sum(a.get("potentialArr", 0) for a in accounts if a.get("play") == play)
        deck.fill(slide, x, BODY_TOP, width, Inches(4.85), PANEL, line=LINE, radius=True)
        deck.fill(slide, x, BODY_TOP, width, Inches(0.06), color)
        deck.text(slide, Emu(int(x + Inches(0.26))), Inches(1.82), Emu(int(width - Inches(0.5))),
                  Inches(0.4), play, size=21, color=color, bold=True)
        deck.text(slide, Emu(int(x + Inches(0.26))), Inches(2.24), Emu(int(width - Inches(0.5))),
                  Inches(0.5), (paf.get("plays", {}).get(play, {}).get("focus", "")),
                  size=11, color=MUTED)
        deck.text(slide, Emu(int(x + Inches(0.26))), Inches(2.92), Emu(int(width - Inches(0.5))),
                  Inches(0.4), "%d accounts" % count, size=17, color=WHITE, bold=True)
        deck.text(slide, Emu(int(x + Inches(0.26))), Inches(3.32), Emu(int(width - Inches(0.5))),
                  Inches(0.4), "%s sized potential" % money(pot), size=12, color=TEXT)
        top = [a for a in accounts if a.get("play") == play][:5]
        deck.text(slide, Emu(int(x + Inches(0.26))), Inches(3.86), Emu(int(width - Inches(0.5))),
                  Inches(0.24), "LEADING ACCOUNTS", size=9, color=MUTED, bold=True, space=True)
        deck.bullets(slide, Emu(int(x + Inches(0.26))), Inches(4.16),
                     Emu(int(width - Inches(0.5))),
                     [truncate(a.get("name", ""), 30) for a in top],
                     size=11, gap=0.1, bullet_color=color)
        land = (paf.get("plays", {}).get(play, {}).get("land") or [])
        if land:
            deck.text(slide, Emu(int(x + Inches(0.26))), Inches(5.62),
                      Emu(int(width - Inches(0.5))), Inches(0.24), "FIRST MOVE",
                      size=9, color=MUTED, bold=True, space=True)
            deck.text(slide, Emu(int(x + Inches(0.26))), Inches(5.92),
                      Emu(int(width - Inches(0.5))), Inches(0.5),
                      truncate(land[0]["title"], 56), size=11.5, color=TEXT, bold=True)
        x += width + Inches(0.25)
    deck.footnote(slide, "Play assignment is deterministic: it is derived from Copilot, security "
                         "and platform signals in the export and does not change between runs.")
    return slide


def slide_potential(deck, potential, focus):
    slide = deck.slide("What the half is worth", "Q3 \u00b7 Potential",
                       note="Every line is tagged observed, list or derived. Observed means the "
                            "account's own billed rate from Kusto. Derived appears only for GHE, "
                            "which has no public per-seat list price.")
    totals = potential.get("totals", {})
    accounts = focus.get("accounts", [])

    seats, committers, ghe_seats = 0, 0, 0
    focus_totals = {"Copilot": 0.0, "GHAS": 0.0, "GHE": 0.0}
    rates = {"Copilot": [], "GHAS": [], "GHE": []}
    bases = {}
    for account in accounts:
        for line in account.get("lines", []):
            product = line.get("product")
            if product in focus_totals:
                focus_totals[product] += float(line.get("value") or 0)
                rates[product].append((float(line.get("rate") or 0), line.get("basis")))
            bases[line.get("basis")] = bases.get(line.get("basis"), 0) + 1
            if product == "Copilot":
                seats += int(line.get("quantity") or 0)
            elif product == "GHAS":
                committers += int(line.get("quantity") or 0)
            elif product == "GHE":
                ghe_seats += int(line.get("quantity") or 0)

    def rate_basis(product):
        """Describe the rate actually applied, rather than restating list price."""
        entries = rates.get(product) or []
        if not entries:
            return "\u2014"
        values = sorted(r for r, _ in entries)
        median = values[len(values) // 2]
        kinds = {}
        for _, basis in entries:
            kinds[basis] = kinds.get(basis, 0) + 1
        modal = max(kinds, key=kinds.get)
        return "%s \u00b7 $%.2f/%s/mo" % (modal, median / 12.0,
                                          "committer" if product == "GHAS" else "user")

    def book(product):
        return float((totals.get(product) or {}).get("value") or 0)

    cards = [
        ("Focus potential ARR", money(sum(focus_totals.values())), "new business only", ACCENT),
        ("Copilot seats", num(seats), "%s in focus set" % money(focus_totals["Copilot"]),
         PLAY_COLOR["Innovate"]),
        ("GHAS committers", num(committers), "%s in focus set" % money(focus_totals["GHAS"]),
         PLAY_COLOR["Trust"]),
        ("GHE seats", num(ghe_seats), "%s in focus set" % money(focus_totals["GHE"]),
         PLAY_COLOR["Scale"]),
    ]
    x = MARGIN
    for label, value, sub, color in cards:
        deck.card(slide, x, BODY_TOP, Inches(2.95), Inches(1.42), label, value, sub, color)
        x += Inches(3.11)

    rows = [
        ["Copilot", "Seats without Copilot today", num(seats), money(focus_totals["Copilot"]),
         money(book("Copilot")), rate_basis("Copilot")],
        ["GHAS", "Active committers L90d not covered", num(committers), money(focus_totals["GHAS"]),
         money(book("GHAS")), rate_basis("GHAS")],
        ["GHE", "Azure DevOps TAM available to migrate", num(ghe_seats), money(focus_totals["GHE"]),
         money(book("GHE")), rate_basis("GHE")],
    ]
    colors = {0: PLAY_COLOR["Innovate"], 1: PLAY_COLOR["Trust"], 2: PLAY_COLOR["Scale"]}
    deck.table(slide, MARGIN, Inches(3.42), W - 2 * MARGIN,
               ["Product", "Sizing basis", "Quantity", "Focus ARR", "Whole book", "Rate basis"],
               [1.2, 4.0, 1.2, 1.4, 1.4, 2.4], rows, row_h=0.42, size=11, colors=colors)

    deck.text(slide, MARGIN, Inches(5.1), W - 2 * MARGIN, Inches(0.3),
              "HOW EACH FIGURE IS PRICED", size=10, color=MUTED, bold=True, space=True)
    deck.bullets(slide, MARGIN, Inches(5.42), W - 2 * MARGIN, [
        "Observed first: where the account already buys the product, its own billed rate from "
        "Kusto is used (%d of %d priced lines)." % (bases.get("observed", 0), sum(bases.values()) or 1),
        "List second: public GitHub pricing for products the account does not yet buy.",
        "Derived only for GHE, which has no published per-seat price \u2014 the median observed "
        "rate across the book is used and labelled as derived wherever it appears.",
    ], size=11.5, gap=0.14)
    deck.footnote(slide, "GHAS is sized per active committer (L90d) because that is how it bills - "
                         "not per seat. Sizing a seat count here would overstate the number.")
    return slide


def slide_aiu(deck, potential, focus):
    slide = deck.slide("AI credits: run-rate today, capacity unlocked by the plan",
                       "Q3 \u00b7 AIU and consumption",
                       note="Credits already invoiced are existing revenue, not upside, so they "
                            "are deliberately excluded from potential ARR. Included credits ship "
                            "bundled with seats - they are capacity, not revenue. Overage is real "
                            "incremental revenue but cannot be forecast without inventing a "
                            "consumption curve, so it is not modelled.")
    aiu = potential.get("aiuTotals", {})
    installed = potential.get("installed", {})

    cards = [
        ("Accounts consuming AIU", num(aiu.get("accountsConsumingAiu")), "measured, not estimated", ACCENT),
        ("Credits consumed / yr", num(aiu.get("currentAnnualisedCredits")), "annualised run-rate", WHITE),
        ("AIU revenue / yr", money(aiu.get("currentAnnualisedSpend")), "already invoiced", WHITE),
        ("Capacity the plan unlocks", num(aiu.get("includedCreditCapacityFromPlan")),
         "included credits, bundled", GOOD),
    ]
    x = MARGIN
    for label, value, sub, color in cards:
        deck.card(slide, x, BODY_TOP, Inches(2.95), Inches(1.42), label, value, sub, color)
        x += Inches(3.11)

    acr = installed.get("acrAnnualised", {}) or {}
    ranked = sorted(acr.items(), key=lambda kv: -kv[1])[:6]
    rows = [[PRODUCT_LABEL.get(name, name.replace("_", " ").capitalize()), money(value)]
            for name, value in ranked]
    panel_h = Inches(0.84) + Inches(0.3) * (len(rows) + 1)
    deck.fill(slide, MARGIN, Inches(3.4), Inches(5.0), panel_h, PANEL, line=LINE, radius=True)
    deck.text(slide, Inches(0.86), Inches(3.6), Inches(4.5), Inches(0.3),
              "ANNUALISED CONSUMPTION REVENUE, WHOLE BOOK", size=10, color=MUTED,
              bold=True, space=True)
    deck.table(slide, Inches(0.86), Inches(3.94), Inches(4.5),
               ["Product", "Annualised"], [2.6, 1.9], rows, row_h=0.3, size=11)

    deck.panel(slide, Inches(6.2), Inches(3.4), Inches(6.5),
               "Why AIU is not in the potential number", [
        "Credits already invoiced are existing revenue. Counting them as upside would "
        "double-count the book.",
        "Included credits (1,900/mo Business, 3,900/mo Enterprise) ship bundled with the "
        "seat, so they are capacity the plan unlocks \u2014 not new revenue.",
        "Overage above the included pool is genuinely incremental, but forecasting it needs "
        "a consumption curve we do not have. It is upside to the number shown, not in it.",
    ], WARN, size=11.5, gap=0.2, min_h=panel_h)
    deck.footnote(slide, "AIU figures come from the Kusto consumption invoice fact, annualised by "
                         "observed billing months rather than by window length.")
    return slide


def slide_how_i_win(deck, play, focus, paf, report):
    accounts = [a for a in focus.get("accounts", []) if a.get("play") == play]
    entry = paf.get("plays", {}).get(play, {})
    color = PLAY_COLOR.get(play, ACCENT)
    pot = sum(a.get("potentialArr", 0) for a in accounts)
    greenfield = sum(1 for a in accounts if (a.get("current") or {}).get("greenfield"))

    slide = deck.slide("How I win %s" % play, "Q4 \u00b7 Execution \u00b7 %d accounts \u00b7 %s"
                       % (len(accounts), money(pot)),
                       note="Sequence is taken from the GitHub Product Adoption Framework key "
                            "actions for this product. Land applies to accounts with no footprint; "
                            "expand applies where the product is already in use.")
    deck.fill(slide, MARGIN, Inches(1.36), Inches(1.5), Inches(0.035), color)

    deck.panel(slide, MARGIN, BODY_TOP, Inches(6.05),
               "Land \u00b7 %d greenfield accounts" % greenfield,
               ["%s \u2014 %s" % (a["title"], truncate(a["summary"], 96))
                for a in entry.get("land", [])[:4]], color, size=11, gap=0.2)

    deck.panel(slide, Inches(7.0), BODY_TOP, Inches(5.7),
               "Expand \u00b7 %d with existing footprint" % (len(accounts) - greenfield),
               ["%s \u2014 %s" % (a["title"], truncate(a["summary"], 90))
                for a in entry.get("expand", [])[:4]], color, size=11, gap=0.2)

    # The PAF sequence answers "what do I run"; leadership also asks "on whom, first".
    # Naming the opening accounts with their own measured signal keeps it specific.
    first = accounts[:3]
    if first:
        deck.text(slide, MARGIN, Inches(5.44), Inches(11), Inches(0.26),
                  "WHERE THIS LANDS FIRST", size=9.5, color=MUTED, bold=True, space=True)
        rows = []
        for account in first:
            signal = ""
            for line in account.get("evidence") or []:
                if line.lower().startswith(play.lower()):
                    signal = line.split(":", 1)[-1].strip()
                    break
            if not signal:
                gaps = account.get("discoveryGaps") or []
                signal = gaps[0] if gaps else "Discovery required"
            sized = ", ".join(
                "%s %s %s" % (num(item.get("quantity")),
                              PRODUCT_LABEL.get(str(item.get("product", "")).lower(),
                                                str(item.get("product", ""))),
                              item.get("metric", ""))
                for item in (account.get("lines") or [])[:2]) or "\u2014"
            rows.append([truncate(account.get("name", ""), 24),
                         money(account.get("potentialArr")),
                         truncate(sized, 34),
                         truncate(signal, 56)])
        deck.table(slide, MARGIN, Inches(5.7), W - 2 * MARGIN,
                   ["Account", "Potential", "What we would sell", "Measured signal to open on"],
                   [2.4, 1.1, 3.4, 5.5], rows, row_h=0.27, size=9.5,
                   colors={i: color for i in range(len(rows))})
    deck.footnote(slide, "Key actions and supporting resources sourced from the GitHub Product "
                         "Adoption Framework. Resource links are listed in the appendix.")
    return slide


def slide_cadence(deck, focus):
    slide = deck.slide("The operating cadence that makes this happen",
                       "Q4 \u00b7 How I win \u00b7 Cadence",
                       note="The plan fails on execution discipline, not on account selection. "
                            "This is the weekly mechanism.")
    tiers = [
        (tier_label(focus, "Tier 1", "Tier 1"), tier_count(focus, "Tier 1"), "Weekly", GOOD, [
            "Named executive sponsor on both sides within 30 days",
            "Live opportunity or PoC in flight by end of Q1",
            "Joint success plan with measurable adoption milestones",
            "Escalation path agreed with the account team",
        ]),
        (tier_label(focus, "Tier 2", "Tier 2"), tier_count(focus, "Tier 2"), "Bi-weekly", ACCENT, [
            "Discovery completed and play validated against real signals",
            "PAF land sequence started with a named champion",
            "Promote to Tier 1 on a qualified opportunity",
            "Demote if no two-way engagement after two cycles",
        ]),
        (tier_label(focus, "Tier 3", "Tier 3"), tier_count(focus, "Tier 3"), "Monthly", MUTED, [
            "Trigger-activated outreach only \u2014 no manufactured touches",
            "Automated nurture against product signals",
            "Re-rank at mid-half against fresh triggers",
            "Promote on funding, breach, leadership change or expansion",
        ]),
    ]
    x = MARGIN
    width = (W - 2 * MARGIN - Inches(0.5)) / 3
    for name, count, cadence, color, items in tiers:
        deck.fill(slide, x, BODY_TOP, width, Inches(3.28), PANEL, line=LINE, radius=True)
        deck.fill(slide, x, BODY_TOP, width, Inches(0.06), color)
        deck.text(slide, Emu(int(x + Inches(0.26))), Inches(1.84), Emu(int(width - Inches(0.5))),
                  Inches(0.4), "%s \u00b7 %d accounts" % (name, count), size=15,
                  color=color, bold=True)
        deck.text(slide, Emu(int(x + Inches(0.26))), Inches(2.26), Emu(int(width - Inches(0.5))),
                  Inches(0.3), "%s review cadence" % cadence, size=11, color=MUTED)
        deck.bullets(slide, Emu(int(x + Inches(0.26))), Inches(2.72),
                     Emu(int(width - Inches(0.5))), items, size=11, gap=0.18, bullet_color=color)
        x += width + Inches(0.25)

    # Cadence answers "how often"; the milestone strip answers "by when", which is the
    # commitment leadership actually holds the half against.
    deck.text(slide, MARGIN, Inches(5.0), Inches(11), Inches(0.26),
              "HALF-LEVEL MILESTONES", size=9.5, color=MUTED, bold=True, space=True)
    tier1 = tier_count(focus, "Tier 1")
    tier2 = tier_count(focus, "Tier 2")
    stages = [
        ("Q1 \u00b7 Weeks 1\u20134", "Tier 1 discovery complete, sponsor named on %d accounts"
         % tier1, GOOD),
        ("Q1 \u00b7 Weeks 5\u201312", "PAF land sequence live; PoC or pilot in flight on Tier 1",
         ACCENT),
        ("Q2 \u00b7 Weeks 13\u201320", "Tier 2 promoted on qualified opportunity; expand motion "
         "on installed base", PLAY_COLOR["Scale"]),
        ("Q2 \u00b7 Weeks 21\u201326", "Close, re-rank against fresh triggers, hand H2 list to "
         "leadership", WARN),
    ]
    sx = MARGIN
    swidth = (W - 2 * MARGIN - Inches(0.36)) / 4
    for label, detail, color in stages:
        deck.fill(slide, sx, Inches(5.32), swidth, Inches(1.08), PANEL, line=LINE, radius=True)
        deck.fill(slide, sx, Inches(5.32), Inches(0.05), Inches(1.08), color)
        deck.text(slide, Emu(int(sx + Inches(0.22))), Inches(5.42),
                  Emu(int(swidth - Inches(0.4))), Inches(0.26), label, size=10,
                  color=color, bold=True)
        deck.text(slide, Emu(int(sx + Inches(0.22))), Inches(5.7),
                  Emu(int(swidth - Inches(0.4))), Inches(0.62), detail, size=10, color=TEXT)
        sx += swidth + Inches(0.12)
    deck.footnote(slide, "Tiers are assigned by composite rank: top 25% Tier 1, next 35% Tier 2, "
                         "remainder Tier 3. They are re-cut at mid-half.")
    return slide


def slide_partners(deck, focus, partners, report):
    accounts = focus.get("accounts", [])
    by_key = (partners or {}).get("accounts", {}) if partners else {}
    with_partner = [a for a in accounts if by_key.get(a.get("key"), {}).get("partners")]
    msft = [a for a in accounts
            if (by_key.get(a.get("key"), {}).get("microsoft") or {}).get("csp")]

    slide = deck.slide("Microsoft overlap and partner leverage", "Q5 \u00b7 Leverage",
                       note="Coverage is stated honestly. Microsoft AM mapping is sparse in the "
                            "source export, so tpid presence is used as the wider signal and the "
                            "gap is called out as an ask rather than papered over.")
    cards = [
        ("Accounts with a partner", str(len(with_partner)),
         "of %d focus accounts" % len(accounts), ACCENT),
        ("Microsoft-linked", str(len(msft)), "MSFT CSP involvement", PLAY_COLOR["Innovate"]),
        ("Partner-led motions", str(sum(len(by_key.get(a.get("key"), {}).get("partners", []))
                                        for a in accounts)), "relationships mapped", GOOD),
        ("Coverage gap", pct(len(accounts) - len(with_partner), len(accounts)),
         "no mapped partner", WARN),
    ]
    x = MARGIN
    for label, value, sub, color in cards:
        deck.card(slide, x, BODY_TOP, Inches(2.95), Inches(1.42), label, value, sub, color)
        x += Inches(3.11)

    rows = []
    for account in accounts:
        entry = by_key.get(account.get("key"), {})
        pdm = ", ".join(entry.get("pdm") or []) or "\u2014"
        for partner in (entry.get("partners") or [])[:1]:
            rows.append([
                truncate(account.get("name", ""), 22),
                truncate(partner.get("name", "\u2014"), 24),
                truncate(partner.get("involvement", "\u2014"), 19),
                truncate(pdm, 18),
                truncate(partner.get("source", "\u2014"), 11),
                money(account.get("potentialArr")),
            ])
        if len(rows) >= 10:
            break

    if rows:
        deck.table(slide, MARGIN, Inches(3.42), W - 2 * MARGIN,
                   ["Account", "Partner", "Involvement", "GitHub PDM", "Source",
                    "Potential ARR"],
                   [2.5, 2.7, 2.1, 2.1, 1.3, 1.4], rows, row_h=0.3, size=10)
    else:
        deck.fill(slide, MARGIN, Inches(3.42), W - 2 * MARGIN, Inches(2.9), PANEL,
                  line=LINE, radius=True)
        deck.text(slide, Inches(0.86), Inches(3.7), Inches(11.6), Inches(0.4),
                  "No partner relationships are mapped for the focus accounts.",
                  size=15, color=WARN, bold=True)
        deck.bullets(slide, Inches(0.86), Inches(4.2), Inches(11.6), [
            "This is a data gap, not a conclusion \u2014 partner records were not available "
            "for this run rather than confirmed absent.",
            "The ask on the next slide requests partner mapping for the Tier 1 accounts so "
            "co-sell motions can be planned rather than improvised.",
            "Microsoft overlap is carried in the source export as tpid; the MSFT AM field is "
            "sparsely populated and cannot be relied on for routing.",
        ], size=12, gap=0.16, bullet_color=WARN)
    deck.footnote(slide, "Partner data from Salesforce partner relationships. Microsoft overlap "
                         "from the tpid and MSFT AM fields in the source export.")
    return slide


def slide_working(deck, report, potential, focus):
    accounts = focus.get("accounts", [])
    total = potential.get("accountsTotal", 0) or 1
    sized = potential.get("accountsSized", 0)
    with_arr = potential.get("accountsWithArr", 0)
    activity = report.get("activity", {}) or {}
    scored = sum(1 for a in report.get("accounts", [])
                 if float((a.get("activity") or {}).get("score") or 0) > 0)
    two_way = sum(1 for a in report.get("accounts", []) if (a.get("activity") or {}).get("twoWay"))
    triggered = focus.get("triggersFound", 0)

    slide = deck.slide("What's working, and what isn't", "Q6 \u00b7 Honest read",
                       note="Stated as measured facts with the numbers attached, so the asks on "
                            "the next slide are evidenced rather than asserted.")
    deck.panel(slide, MARGIN, BODY_TOP, Inches(6.05), "Working", [
        "Play classification is complete and deterministic \u2014 every one of the %d accounts "
        "carries an evidenced play." % total,
        "%d of %d accounts (%s) now have dollar-sized potential from real product signals "
        "rather than estimates." % (sized, total, pct(sized, total)),
        "%d focus accounts have a dated, cited live trigger giving a genuine reason to call."
        % triggered,
        "%d accounts show verified two-way conversation, which is where conversion is "
        "concentrating." % two_way,
        "Consumption revenue is visible per account from billing facts, so expansion "
        "conversations start from what the customer actually uses.",
    ], GOOD, size=11.5, gap=0.2)

    deck.panel(slide, Inches(7.0), BODY_TOP, Inches(5.7), "Not working", [
        "Only %d of %d accounts carry any ARR. Pipeline creation, not renewal, is the "
        "constraint this half." % (with_arr, total),
        "Activity coverage is thin \u2014 %d accounts scored and %d two-way out of %d. Most of "
        "the book has no recorded conversation at all." % (scored, two_way, total),
        "%d accounts cannot be sized because product signals are missing; they need "
        "discovery before they can be forecast." % (total - sized),
        "Microsoft AM mapping is sparse in the source data, so co-sell routing is manual.",
        "GHE has no public per-seat price, so its sizing uses a derived median and is the "
        "least defensible line in the plan.",
    ], WARN, size=11.5, gap=0.2)
    deck.footnote(slide, "Every figure on this slide is computed from the run, not asserted.")
    return slide


def slide_asks(deck, report, potential, focus):
    total = potential.get("accountsTotal", 0) or 1
    sized = potential.get("accountsSized", 0)
    tier1 = tier_count(focus, "Tier 1")
    slide = deck.slide("What I need from leadership and supporting functions",
                       "Q7 \u00b7 Asks",
                       note="Each ask maps to a specific measured gap on the previous slide.")
    book = potential.get("totals", {}) or {}
    ghas_value = (book.get("GHAS") or {}).get("value") or 0
    asks = [
        ("Leadership", ACCENT, [
            "Executive sponsor assigned to each of the %d Tier 1 accounts" % tier1,
            "Air cover for a concentrated %d-account half rather than full-book coverage"
            % focus.get("selectedCount", 0),
            "Agreement that pipeline creation is the H1 measure, given only %d accounts "
            "carry ARR" % potential.get("accountsWithArr", 0),
            "Deal-desk latitude on GHE pricing where no list rate exists",
        ]),
        ("Partnerships", GOOD, [
            "Partner mapping for the Tier 1 accounts so co-sell is planned, not improvised",
            "Named Microsoft counterpart per tpid-linked account \u2014 the MSFT AM field "
            "cannot be relied on today",
            "GSI and LSP introductions where the play is a platform migration",
            "Joint capacity with the incumbent partners already mapped in this book",
        ]),
        ("Marketing & Enablement", PLAY_COLOR["Innovate"], [
            "Trigger-based campaign support for the Tier 3 nurture list",
            "Copilot and GHAS workshop capacity aligned to the PAF land sequence",
            "Customer references in the same vertical as the Tier 1 accounts",
            "Security-buyer content to carry the %s GHAS opportunity" % money(ghas_value),
        ]),
        ("Data & Operations", WARN, [
            "Fill product-signal gaps for the %d accounts that cannot be sized" % (total - sized),
            "Improve activity capture \u2014 coverage is too thin to rank on reliably",
            "Publish a defensible GHE per-seat rate to replace the derived median",
            "Backfill EA renewal dates \u2014 populated for only a handful of accounts",
        ]),
    ]
    x, y = MARGIN, BODY_TOP
    width = (W - 2 * MARGIN - Inches(0.28)) / 2
    for index, (owner, color, items) in enumerate(asks):
        cx = x if index % 2 == 0 else x + width + Inches(0.28)
        cy = y if index < 2 else y + Inches(2.68)
        deck.fill(slide, cx, cy, width, Inches(2.5), PANEL, line=LINE, radius=True)
        deck.fill(slide, cx, cy, Inches(0.05), Inches(2.5), color)
        deck.text(slide, Emu(int(cx + Inches(0.26))), Emu(int(cy + Inches(0.2))),
                  Emu(int(width - Inches(0.5))), Inches(0.32), owner.upper(),
                  size=11, color=color, bold=True, space=True)
        deck.bullets(slide, Emu(int(cx + Inches(0.26))), Emu(int(cy + Inches(0.64))),
                     Emu(int(width - Inches(0.5))), items, size=11, gap=0.16, bullet_color=color)
    deck.footnote(slide, "Asks are derived from the measured gaps on the previous slide.")
    return slide


def slide_method(deck, report, potential, focus):
    slide = deck.slide("How this was built", "Appendix \u00b7 Methodology",
                       note="Included so any number can be challenged and traced.")
    weights = focus.get("weights", {})
    deck.panel(slide, MARGIN, BODY_TOP, Inches(6.05), "Selection", [
        "Composite score = %.0f%% potential ARR, %.0f%% active communication, %.0f%% live "
        "trigger." % (100 * weights.get("potential", 0), 100 * weights.get("communication", 0),
                      100 * weights.get("trigger", 0)),
        "Potential is log-compressed before scoring because it is heavily skewed \u2014 without "
        "it the top few accounts would flatten everything else to zero.",
        "Ranking runs in two stages so triggers are only researched for accounts that could "
        "realistically make the list.",
        "Triggers decay with a 120-day half-life. Any trigger without a date and a source URL "
        "is dropped, not discounted.",
        "Tiers: top 25% Tier 1, next 35% Tier 2, remainder Tier 3.",
    ], ACCENT, size=11, gap=0.2)

    deck.panel(slide, Inches(7.0), BODY_TOP, Inches(5.7), "Sources", [
        "Account list, product signals and plays: %s"
        % truncate(report.get("sourceName", "SuperDash export"), 44),
        "Current ARR and seats: Kusto salesforce_account_monthly_product_arr_seats_fact.",
        "Consumption, ACR and AI credits: Kusto salesforce_account_consumption_invoices_fact.",
        "Communication scoring: Salesforce activity (tasks, events, meetings).",
        "Execution guidance: GitHub Product Adoption Framework key actions.",
        "Live triggers: public web sources, each carrying a date and a URL.",
    ], ACCENT, size=11, gap=0.2)
    deck.footnote(slide, "Generated %s. Re-running the pipeline on the same inputs reproduces "
                         "this deck exactly." % date.today().isoformat())
    return slide


def slide_appendix(deck, focus, index, total, chunk):
    slide = deck.slide("Focus account detail (%d of %d)" % (index, total),
                       "Appendix \u00b7 Account detail")
    rows, colors = [], {}
    for i, account in enumerate(chunk):
        triggers = account.get("triggers") or []
        rows.append([
            truncate(account.get("name", ""), 24),
            account.get("tier", "\u2014"),
            account.get("play", "\u2014"),
            money(account.get("potentialArr")),
            money((account.get("current") or {}).get("arr")),
            "%.0f" % float(account.get("compositeScore") or 0),
            str(len(triggers)),
            truncate(account.get("renewal") or "\u2014", 10),
        ])
        colors[i] = PLAY_COLOR.get(account.get("play"), TEXT)
    deck.table(slide, MARGIN, Inches(1.7), W - 2 * MARGIN,
               ["Account", "Tier", "Play", "Potential ARR", "Current ARR", "Score",
                "Triggers", "Renewal"],
               [3.2, 1.0, 1.2, 1.5, 1.4, 0.9, 1.0, 1.3], rows, row_h=0.275, size=9.5,
               colors=colors)
    deck.footnote(slide, "Score is the composite of potential, communication and trigger recency.")
    return slide


def slide_triggers(deck, focus):
    accounts = [a for a in focus.get("accounts", []) if a.get("triggers")]
    if not accounts:
        return None
    slide = deck.slide("Live triggers behind the ranking", "Appendix \u00b7 Trigger evidence",
                       note="Every trigger carries a date and a source URL. Undated or uncited "
                            "claims were dropped during collection rather than discounted.")
    rows, colors = [], {}
    flat = []
    for account in accounts:
        for trigger in account.get("triggers", []):
            flat.append((account, trigger))
    for i, (account, trigger) in enumerate(flat[:16]):
        first = i == 0 or flat[i - 1][0] is not account
        rows.append([
            truncate(account.get("name", ""), 24) if first else "",
            truncate(trigger.get("type", "\u2014"), 18),
            truncate(trigger.get("date", "\u2014"), 11),
            truncate(trigger.get("headline", ""), 62),
        ])
        colors[i] = PLAY_COLOR.get(account.get("play"), TEXT)
    deck.table(slide, MARGIN, Inches(1.7), W - 2 * MARGIN,
               ["Account", "Trigger type", "Date", "Headline"],
               [3.0, 2.0, 1.4, 5.7], rows, row_h=0.3, size=9.5, colors=colors)
    deck.footnote(slide, "Full source URLs are carried in focus-accounts.json and the workbook.")
    return slide


def slide_paf_resources(deck, paf):
    slide = deck.slide("Execution resources by play", "Appendix \u00b7 PAF resources",
                       note="Resources come from the GitHub Product Adoption Framework key "
                            "action mapping.")
    x = MARGIN
    width = (W - 2 * MARGIN - Inches(0.5)) / 3
    for play in ("Innovate", "Trust", "Scale"):
        color = PLAY_COLOR.get(play, ACCENT)
        entry = paf.get("plays", {}).get(play, {})
        deck.fill(slide, x, BODY_TOP, width, Inches(4.85), PANEL, line=LINE, radius=True)
        deck.fill(slide, x, BODY_TOP, width, Inches(0.06), color)
        deck.text(slide, Emu(int(x + Inches(0.26))), Inches(1.84), Emu(int(width - Inches(0.5))),
                  Inches(0.36), play, size=16, color=color, bold=True)
        items = []
        for action in (entry.get("land", []) + entry.get("expand", []))[:6]:
            for resource in action.get("resources", [])[:1]:
                items.append("%s \u2014 %s" % (truncate(resource["name"], 40), resource["provider"]))
        deck.bullets(slide, Emu(int(x + Inches(0.26))), Inches(2.32),
                     Emu(int(width - Inches(0.5))), items, size=10.5, gap=0.2,
                     bullet_color=color)
        x += width + Inches(0.25)
    deck.footnote(slide, "Links are carried in paf.json alongside the deck for direct sharing.")
    return slide


# ---------------------------------------------------------------- entry point

def load(path, default=None):
    if not path or not os.path.exists(path):
        return default
    with open(path, encoding="utf-8") as fh:
        return json.load(fh)


def main():
    args = sys.argv[1:]
    if len(args) < 4:
        raise SystemExit(
            "usage: deck.py <report.json> <potential.json> <focus-accounts.json> <runDir> "
            "[--paf f] [--partners f] [--out f]")

    report_path, potential_path, focus_path, run_dir = args[:4]

    def opt(flag, default=None):
        return args[args.index(flag) + 1] if flag in args else default

    here = os.path.dirname(os.path.abspath(__file__))
    report = load(report_path) or {}
    potential = load(potential_path) or {}
    focus = load(focus_path) or {}
    paf = load(opt("--paf", os.path.join(here, "..", "paf.json")), {"plays": {}})
    partners = load(opt("--partners", os.path.join(run_dir, "partners.json")))
    out_path = opt("--out", os.path.join(run_dir, "fy27-h1-focus-accounts.pptx"))

    if not focus.get("accounts"):
        raise SystemExit("focus-accounts.json has no accounts - run rank.py stage2 first.")

    deck = Deck()
    slide_title(deck, report, focus, potential)
    slide_thesis(deck, report, focus, potential)

    tiers = tier_names(focus)
    for index, tier in enumerate(tiers, 1):
        slide_accounts(deck, focus, tier, index, len(tiers))

    slide_plays(deck, report, focus, paf)
    slide_potential(deck, potential, focus)
    slide_aiu(deck, potential, focus)
    for play in ("Innovate", "Trust", "Scale"):
        if any(a.get("play") == play for a in focus["accounts"]):
            slide_how_i_win(deck, play, focus, paf, report)
    slide_cadence(deck, focus)
    slide_partners(deck, focus, partners, report)
    slide_working(deck, report, potential, focus)
    slide_asks(deck, report, potential, focus)
    slide_method(deck, report, potential, focus)

    accounts = focus["accounts"]
    chunks = [accounts[i:i + 18] for i in range(0, len(accounts), 18)]
    for index, chunk in enumerate(chunks, 1):
        slide_appendix(deck, focus, index, len(chunks), chunk)
    slide_triggers(deck, focus)
    slide_paf_resources(deck, paf)

    deck.save(out_path)
    print(json.dumps({
        "deckPath": os.path.abspath(out_path),
        "slides": len(deck.prs.slides.__iter__.__self__._sldIdLst),
        "focusAccounts": len(accounts),
        "tiers": focus.get("tierMix", {}),
        "potentialArr": focus.get("totals", {}).get("potentialArr"),
    }))
    return 0


if __name__ == "__main__":
    sys.exit(main())
